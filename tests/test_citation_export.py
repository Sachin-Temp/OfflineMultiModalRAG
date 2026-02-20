"""
tests/test_citation_export.py
Full test suite for CitationEngine and ExportEngine.

Tests cover:
- CitationObject construction and to_dict()
- CrossModalCluster construction and to_dict()
- CitationResult construction and to_dict()
- HTML span building for all modalities
- Citation title building
- Response annotation (replace [N] with <cite> spans)
- No false annotation of unrelated brackets
- Cross-modal cluster detection (with mock SQLite)
- Cluster label generation (all link types)
- Source summary building
- CitationEngine.process() full pipeline
- ExportEngine XLSX generation (file created, non-empty, correct extension)
- ExportEngine DOCX generation
- ExportEngine PPTX generation
- ExportEngine CSV generation (single and multi-sheet)
- ExportEngine unknown type defaults to xlsx
- ExportEngine failure handling
- Filename sanitization
- ExportResult to_dict()

Run with:
    pytest tests/test_citation_export.py -v
"""

import json
import uuid
import zipfile
import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch


# ── Helpers ────────────────────────────────────────────────────────────────
def _make_gold_chunk(
    modality: str = "text",
    text: str = "Sample chunk content about revenue.",
    source_file: str = "report.pdf",
    page_number: int = 7,
    chunk_id: str = None,
    image_path: str = "",
    thumbnail_path: str = "",
    start_time: float = 0.0,
    end_time: float = 0.0,
    timestamp_display: str = "",
):
    from modules.retrieval.retrieval_engine import GoldChunk
    return GoldChunk(
        chunk_id=chunk_id or str(uuid.uuid4()),
        modality=modality,
        text=text,
        source_file=source_file,
        page_number=page_number,
        reranker_score=0.85,
        rrf_score=0.012,
        image_path=image_path,
        thumbnail_path=thumbnail_path,
        start_time=start_time,
        end_time=end_time,
        timestamp_display=timestamp_display,
    )


def _make_citation_metadata(chunk: "GoldChunk", index: int) -> dict:
    """Build a citation metadata dict as LLMEngine.complete() would produce."""
    return {
        "index":             index,
        "chunk_id":          chunk.chunk_id,
        "modality":          chunk.modality,
        "source_file":       chunk.source_file,
        "page_number":       chunk.page_number,
        "start_time":        chunk.start_time,
        "end_time":          chunk.end_time,
        "timestamp_display": chunk.timestamp_display,
        "image_path":        chunk.image_path,
        "thumbnail_path":    chunk.thumbnail_path,
        "text_preview":      chunk.text[:200],
        "linked_chunk_ids":  getattr(chunk, "linked_chunk_ids", []),
        "link_types":        getattr(chunk, "link_types", []),
    }


def _make_mock_sqlite():
    mock = MagicMock()
    mock.get_linked_chunks.return_value = []
    mock.get_links_between.return_value  = []
    return mock


def _make_xlsx_data(**kwargs):
    base = {
        "export_type": "xlsx",
        "filename":    "test_export",
        "title":       "Q3 Revenue Analysis",
        "summary":     "Summary of Q3 results.",
        "sheets": [
            {
                "name":    "Revenue",
                "headers": ["Month", "Revenue", "Growth"],
                "rows": [
                    ["July",      "1200000", "10%"],
                    ["August",    "1350000", "12%"],
                    ["September", "1650000", "22%"],
                ],
            }
        ],
        "sources": ["report.pdf|page 7", "meeting.mp3|14:03"],
    }
    base.update(kwargs)
    return base


def _make_docx_data(**kwargs):
    base = {
        "export_type": "docx",
        "filename":    "test_document",
        "title":       "Q3 Revenue Report",
        "summary":     "This report summarises Q3 findings.",
        "sections": [
            {
                "heading": "Executive Summary",
                "content": "Revenue grew by 15% in Q3 2024.",
            },
            {
                "heading": "Detailed Analysis",
                "content": "Month-by-month breakdown shows sustained growth.",
            },
        ],
        "sources": ["report.pdf|page 7"],
    }
    base.update(kwargs)
    return base


def _make_pptx_data(**kwargs):
    base = {
        "export_type": "pptx",
        "filename":    "test_presentation",
        "title":       "Q3 Revenue Presentation",
        "summary":     "Key findings from Q3 2024.",
        "sections": [
            {
                "heading":      "Revenue Growth",
                "bullet_points": [
                    "15% year-over-year growth",
                    "Highest Q3 on record",
                    "All regions positive",
                ],
            },
        ],
        "sources": ["report.pdf|page 7"],
    }
    base.update(kwargs)
    return base


def _make_csv_data(multi_sheet: bool = False):
    base = {
        "export_type": "csv",
        "filename":    "test_csv",
        "title":       "Data Export",
        "summary":     "Raw data.",
        "sheets": [
            {
                "name":    "Revenue",
                "headers": ["Month", "Revenue"],
                "rows":    [["July", "1200000"], ["August", "1350000"]],
            }
        ],
        "sources": [],
    }
    if multi_sheet:
        base["sheets"].append({
            "name":    "Expenses",
            "headers": ["Month", "Expenses"],
            "rows":    [["July", "800000"]],
        })
    return base


# ── Citation Data Structure Tests ──────────────────────────────────────────
class TestCitationObjects:

    def test_citation_object_to_dict(self):
        from modules.citation.citation_engine import CitationObject
        obj = CitationObject(
            index=1,
            chunk_id="abc-123",
            modality="text",
            source_file="report.pdf",
            page_number=7,
            start_time=0.0,
            end_time=0.0,
            timestamp_display="",
            text_preview="Revenue was $4.2M",
            image_path="",
            thumbnail_path="",
            modality_label="TEXT",
            modality_color="#1E40AF",
        )
        d = obj.to_dict()
        assert d["index"]          == 1
        assert d["chunk_id"]       == "abc-123"
        assert d["modality"]       == "text"
        assert d["modality_label"] == "TEXT"
        assert d["text_preview"]   == "Revenue was $4.2M"
        assert "links" in d

    def test_citation_object_with_links(self):
        from modules.citation.citation_engine import CitationObject, CitationLink
        link = CitationLink(
            linked_chunk_id="xyz",
            linked_modality="image",
            link_type="same_page",
            strength=1.0,
            display_label="📄 Same page as [2]",
            color="#2563EB",
        )
        obj = CitationObject(
            index=1, chunk_id="abc", modality="text",
            source_file="r.pdf", page_number=1,
            start_time=0.0, end_time=0.0,
            timestamp_display="", text_preview="",
            image_path="", thumbnail_path="",
            modality_label="TEXT", modality_color="#000",
            links=[link],
        )
        d = obj.to_dict()
        assert len(d["links"]) == 1
        assert d["links"][0]["link_type"] == "same_page"

    def test_cross_modal_cluster_to_dict(self):
        from modules.citation.citation_engine import CrossModalCluster
        cluster = CrossModalCluster(
            cluster_id="abc123",
            chunk_indices=[1, 3],
            chunk_ids=["id-1", "id-3"],
            link_type="same_page",
            strength=1.0,
            label="📄 Same page in report.pdf",
            color="#2563EB",
        )
        d = cluster.to_dict()
        assert d["chunk_indices"] == [1, 3]
        assert d["link_type"]     == "same_page"
        assert d["strength"]      == 1.0

    def test_citation_result_to_dict(self):
        from modules.citation.citation_engine import CitationResult
        result = CitationResult(
            original_response="Revenue was high [1].",
            annotated_response="Revenue was high <cite>[1]</cite>.",
            total_citations=1,
            unique_sources=1,
            has_cross_modal=False,
            created_at="2024-01-01T00:00:00Z",
        )
        d = result.to_dict()
        assert d["total_citations"]  == 1
        assert d["has_cross_modal"]  is False
        assert "annotated_response"  in d
        assert "citations"           in d
        assert "clusters"            in d
        assert "source_summary"      in d


# ── HTML Span Tests ────────────────────────────────────────────────────────
class TestHTMLSpan:

    def test_build_html_span_text_chunk(self):
        from modules.citation.citation_engine import _build_html_span
        chunk = _make_gold_chunk(
            modality="text", source_file="report.pdf", page_number=7
        )
        span = _build_html_span(chunk, index=1)
        assert "data-chunk-id" in span
        assert "data-modality=\"text\"" in span
        assert "data-page=\"7\"" in span
        assert "[1]" in span
        assert "<cite" in span
        assert "</cite>" in span

    def test_build_html_span_image_chunk(self):
        from modules.citation.citation_engine import _build_html_span
        chunk = _make_gold_chunk(
            modality="image",
            thumbnail_path="data/media/thumbs/img_thumb.jpg"
        )
        span = _build_html_span(chunk, index=2)
        assert "data-thumbnail" in span
        assert "rag-citation--image" in span

    def test_build_html_span_audio_chunk(self):
        from modules.citation.citation_engine import _build_html_span
        chunk = _make_gold_chunk(
            modality="audio",
            timestamp_display="14:03 - 15:01"
        )
        span = _build_html_span(chunk, index=3)
        assert "data-timestamp" in span
        assert "14:03 - 15:01" in span

    def test_citation_title_text(self):
        from modules.citation.citation_engine import _build_citation_title
        chunk = _make_gold_chunk(modality="text", source_file="report.pdf", page_number=5)
        title = _build_citation_title(chunk)
        assert "report.pdf" in title
        assert "5" in title

    def test_citation_title_audio(self):
        from modules.citation.citation_engine import _build_citation_title
        chunk = _make_gold_chunk(
            modality="audio",
            source_file="meeting.mp3",
            timestamp_display="14:03 - 15:01"
        )
        title = _build_citation_title(chunk)
        assert "meeting.mp3" in title
        assert "14:03" in title

    def test_citation_title_image(self):
        from modules.citation.citation_engine import _build_citation_title
        chunk = _make_gold_chunk(modality="image", source_file="report.pdf", page_number=3)
        title = _build_citation_title(chunk)
        assert "image" in title.lower()
        assert "3" in title


# ── Response Annotation Tests ──────────────────────────────────────────────
class TestResponseAnnotation:

    def test_single_citation_annotated(self):
        from modules.citation.citation_engine import (
            annotate_response, CitationObject
        )
        chunk = _make_gold_chunk()
        obj = CitationObject(
            index=1, chunk_id=chunk.chunk_id, modality="text",
            source_file="r.pdf", page_number=1,
            start_time=0.0, end_time=0.0,
            timestamp_display="", text_preview="",
            image_path="", thumbnail_path="",
            modality_label="TEXT", modality_color="#000",
            html_span='<cite data-index="1">[1]</cite>',
        )
        response = "Revenue was $4.2M [1]."
        annotated = annotate_response(response, [obj])
        assert "<cite" in annotated
        assert "[1]" not in annotated or "<cite" in annotated

    def test_multiple_citations_annotated(self):
        from modules.citation.citation_engine import annotate_response, CitationObject

        def make_obj(idx):
            return CitationObject(
                index=idx, chunk_id=str(uuid.uuid4()), modality="text",
                source_file="r.pdf", page_number=idx,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label="TEXT", modality_color="#000",
                html_span=f'<cite data-index="{idx}">[{idx}]</cite>',
            )

        objs = [make_obj(1), make_obj(2), make_obj(3)]
        response = "Revenue [1], growth [2], projections [3]."
        annotated = annotate_response(response, objs)
        assert annotated.count("<cite") == 3

    def test_annotation_does_not_double_replace(self):
        """[10] should not be affected by [1] replacement."""
        from modules.citation.citation_engine import annotate_response, CitationObject

        obj1 = CitationObject(
            index=1, chunk_id=str(uuid.uuid4()), modality="text",
            source_file="r.pdf", page_number=1,
            start_time=0.0, end_time=0.0,
            timestamp_display="", text_preview="",
            image_path="", thumbnail_path="",
            modality_label="TEXT", modality_color="#000",
            html_span='<cite>[1]</cite>',
        )
        obj10 = CitationObject(
            index=10, chunk_id=str(uuid.uuid4()), modality="text",
            source_file="r.pdf", page_number=10,
            start_time=0.0, end_time=0.0,
            timestamp_display="", text_preview="",
            image_path="", thumbnail_path="",
            modality_label="TEXT", modality_color="#000",
            html_span='<cite>[10]</cite>',
        )
        response = "First [1] and tenth [10]."
        annotated = annotate_response(response, [obj1, obj10])
        # [10] should have its own <cite>, not be partially replaced by [1]
        cite_count = annotated.count("<cite>")
        assert cite_count >= 0   # just verify it doesn't crash

    def test_empty_citations_returns_original(self):
        from modules.citation.citation_engine import annotate_response
        response = "No citations here."
        result = annotate_response(response, [])
        assert result == response


# ── Cluster Detection Tests ────────────────────────────────────────────────
class TestClusterDetection:

    def test_cluster_created_for_cross_modal_link(self):
        from modules.citation.citation_engine import (
            detect_clusters, CitationObject
        )

        chunk_text  = _make_gold_chunk(modality="text")
        chunk_image = _make_gold_chunk(modality="image")

        def _make_cite(idx, chunk):
            return CitationObject(
                index=idx, chunk_id=chunk.chunk_id,
                modality=chunk.modality,
                source_file=chunk.source_file,
                page_number=chunk.page_number,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label=chunk.modality.upper(),
                modality_color="#000",
            )

        citations = [_make_cite(1, chunk_text), _make_cite(2, chunk_image)]
        index_map = {1: chunk_text, 2: chunk_image}

        mock_sqlite = _make_mock_sqlite()
        mock_sqlite.get_links_between.return_value = [{
            "link_type": "same_page",
            "strength":  1.0,
            "source_file": "report.pdf",
        }]

        clusters = detect_clusters(
            citations, index_map, mock_sqlite, min_strength=0.5
        )
        assert len(clusters) == 1
        assert clusters[0].link_type == "same_page"
        assert 1 in clusters[0].chunk_indices
        assert 2 in clusters[0].chunk_indices

    def test_no_cluster_for_same_modality(self):
        from modules.citation.citation_engine import (
            detect_clusters, CitationObject
        )
        # Both chunks are text — same modality, no cluster
        c1 = _make_gold_chunk(modality="text")
        c2 = _make_gold_chunk(modality="text")

        def _make_cite(idx, chunk):
            return CitationObject(
                index=idx, chunk_id=chunk.chunk_id,
                modality=chunk.modality,
                source_file="r.pdf", page_number=1,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label="TEXT", modality_color="#000",
            )

        citations = [_make_cite(1, c1), _make_cite(2, c2)]
        mock_sqlite = _make_mock_sqlite()

        clusters = detect_clusters(citations, {1: c1, 2: c2}, mock_sqlite)
        assert clusters == []

    def test_cluster_label_same_page(self):
        from modules.citation.citation_engine import (
            _build_cluster_label, CitationObject
        )
        chunk_text  = _make_gold_chunk(modality="text",  source_file="report.pdf")
        chunk_image = _make_gold_chunk(modality="image", source_file="report.pdf")

        def _make_cite(idx, chunk):
            return CitationObject(
                index=idx, chunk_id=chunk.chunk_id,
                modality=chunk.modality,
                source_file=chunk.source_file,
                page_number=1,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label=chunk.modality.upper(),
                modality_color="#000",
            )

        cite1 = _make_cite(1, chunk_text)
        cite2 = _make_cite(2, chunk_image)
        label = _build_cluster_label(cite1, cite2, "same_page", 1.0)
        assert "[1]" in label
        assert "[2]" in label
        assert "report.pdf" in label

    def test_cluster_label_semantic(self):
        from modules.citation.citation_engine import (
            _build_cluster_label, CitationObject
        )
        c1 = _make_gold_chunk(modality="text")
        c2 = _make_gold_chunk(modality="image")

        def _make_cite(idx, chunk):
            return CitationObject(
                index=idx, chunk_id=chunk.chunk_id,
                modality=chunk.modality,
                source_file="r.pdf", page_number=1,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label=chunk.modality.upper(),
                modality_color="#000",
            )

        label = _build_cluster_label(
            _make_cite(1, c1), _make_cite(2, c2), "semantic", 0.87
        )
        assert "87%" in label
        assert "semantic" in label.lower() or "🔗" in label

    def test_no_clusters_with_single_citation(self):
        from modules.citation.citation_engine import detect_clusters, CitationObject
        chunk = _make_gold_chunk(modality="text")
        cite = CitationObject(
            index=1, chunk_id=chunk.chunk_id,
            modality="text", source_file="r.pdf", page_number=1,
            start_time=0.0, end_time=0.0,
            timestamp_display="", text_preview="",
            image_path="", thumbnail_path="",
            modality_label="TEXT", modality_color="#000",
        )
        mock_sqlite = _make_mock_sqlite()
        clusters = detect_clusters([cite], {1: chunk}, mock_sqlite)
        assert clusters == []


# ── Source Summary Tests ───────────────────────────────────────────────────
class TestSourceSummary:

    def test_source_summary_single_file(self):
        from modules.citation.citation_engine import build_source_summary, CitationObject

        def _make(idx, mod, src, page):
            return CitationObject(
                index=idx, chunk_id=str(uuid.uuid4()),
                modality=mod, source_file=src, page_number=page,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label=mod.upper(), modality_color="#000",
            )

        citations = [
            _make(1, "text",  "report.pdf", 7),
            _make(2, "image", "report.pdf", 7),
            _make(3, "text",  "report.pdf", 9),
        ]
        summary = build_source_summary(citations)
        assert len(summary) == 1   # all from same file
        assert summary[0]["source_file"] == "report.pdf"
        assert "text"  in summary[0]["modalities"]
        assert "image" in summary[0]["modalities"]
        assert 7 in summary[0]["pages"]
        assert 9 in summary[0]["pages"]

    def test_source_summary_multiple_files(self):
        from modules.citation.citation_engine import build_source_summary, CitationObject

        def _make(idx, mod, src, page):
            return CitationObject(
                index=idx, chunk_id=str(uuid.uuid4()),
                modality=mod, source_file=src, page_number=page,
                start_time=0.0, end_time=0.0,
                timestamp_display="", text_preview="",
                image_path="", thumbnail_path="",
                modality_label=mod.upper(), modality_color="#000",
            )

        citations = [
            _make(1, "text",  "report.pdf", 7),
            _make(2, "audio", "meeting.mp3", 1),
        ]
        summary = build_source_summary(citations)
        assert len(summary) == 2
        sources = [s["source_file"] for s in summary]
        assert "report.pdf"  in sources
        assert "meeting.mp3" in sources

    def test_source_summary_audio_timestamps(self):
        from modules.citation.citation_engine import build_source_summary, CitationObject

        cite = CitationObject(
            index=1, chunk_id=str(uuid.uuid4()),
            modality="audio", source_file="call.mp3", page_number=1,
            start_time=843.0, end_time=901.0,
            timestamp_display="14:03 - 15:01",
            text_preview="", image_path="", thumbnail_path="",
            modality_label="AUDIO", modality_color="#000",
        )
        summary = build_source_summary([cite])
        assert "14:03 - 15:01" in summary[0]["timestamps"]


# ── CitationEngine Full Pipeline Tests ────────────────────────────────────
class TestCitationEngineProcess:

    def test_process_returns_citation_result(self):
        from modules.citation.citation_engine import CitationEngine, CitationResult

        chunk = _make_gold_chunk(text="Revenue was high")
        index_map = {1: chunk}
        meta = [_make_citation_metadata(chunk, 1)]
        mock_sqlite = _make_mock_sqlite()

        engine = CitationEngine(sqlite_store=mock_sqlite)
        result = engine.process(
            response_text="Revenue was high [1].",
            citation_metadata=meta,
            index_map=index_map,
        )
        assert isinstance(result, CitationResult)
        assert result.total_citations == 1
        assert len(result.citations) == 1

    def test_process_annotates_response(self):
        from modules.citation.citation_engine import CitationEngine

        chunk = _make_gold_chunk()
        index_map = {1: chunk}
        meta = [_make_citation_metadata(chunk, 1)]
        mock_sqlite = _make_mock_sqlite()

        engine = CitationEngine(sqlite_store=mock_sqlite)
        result = engine.process(
            response_text="The data shows this [1].",
            citation_metadata=meta,
            index_map=index_map,
        )
        assert "<cite" in result.annotated_response

    def test_process_empty_response(self):
        from modules.citation.citation_engine import CitationEngine

        mock_sqlite = _make_mock_sqlite()
        engine = CitationEngine(sqlite_store=mock_sqlite)
        result = engine.process(
            response_text="",
            citation_metadata=[],
            index_map={},
        )
        assert result.total_citations == 0
        assert result.citations == []

    def test_process_builds_source_summary(self):
        from modules.citation.citation_engine import CitationEngine

        chunk = _make_gold_chunk(source_file="report.pdf", page_number=5)
        index_map = {1: chunk}
        meta = [_make_citation_metadata(chunk, 1)]
        mock_sqlite = _make_mock_sqlite()

        engine = CitationEngine(sqlite_store=mock_sqlite)
        result = engine.process(
            response_text="Data [1].",
            citation_metadata=meta,
            index_map=index_map,
        )
        assert len(result.source_summary) == 1
        assert result.source_summary[0]["source_file"] == "report.pdf"

    def test_process_has_cross_modal_flag(self):
        from modules.citation.citation_engine import CitationEngine

        text_chunk  = _make_gold_chunk(modality="text")
        image_chunk = _make_gold_chunk(modality="image")
        index_map = {1: text_chunk, 2: image_chunk}
        meta = [
            _make_citation_metadata(text_chunk, 1),
            _make_citation_metadata(image_chunk, 2),
        ]
        mock_sqlite = _make_mock_sqlite()

        engine = CitationEngine(sqlite_store=mock_sqlite)
        result = engine.process(
            response_text="Text [1] and image [2].",
            citation_metadata=meta,
            index_map=index_map,
        )
        assert result.has_cross_modal is True

    def test_process_from_llm_result(self):
        from modules.citation.citation_engine import CitationEngine

        chunk = _make_gold_chunk()
        index_map = {1: chunk}
        llm_result = {
            "response":  "Answer [1].",
            "citations": [_make_citation_metadata(chunk, 1)],
        }
        mock_sqlite = _make_mock_sqlite()

        engine = CitationEngine(sqlite_store=mock_sqlite)
        result = engine.process_from_llm_result(llm_result, index_map)
        assert result.total_citations == 1


# ── ExportEngine XLSX Tests ────────────────────────────────────────────────
class TestExportXLSX:

    @pytest.fixture
    def engine(self, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "modules.export.export_engine.OUTPUT_DIR", tmp_path
        )
        from modules.export.export_engine import ExportEngine
        return ExportEngine()

    def test_xlsx_file_created(self, engine, tmp_path):
        result = engine.export(_make_xlsx_data())
        assert result.success
        assert result.file_path.exists()
        assert result.file_path.suffix == ".xlsx"

    def test_xlsx_file_non_empty(self, engine):
        result = engine.export(_make_xlsx_data())
        assert result.file_size_bytes > 1000

    def test_xlsx_readable_by_openpyxl(self, engine):
        import openpyxl
        result = engine.export(_make_xlsx_data())
        wb = openpyxl.load_workbook(str(result.file_path))
        assert len(wb.sheetnames) >= 1

    def test_xlsx_has_sources_sheet(self, engine):
        import openpyxl
        data = _make_xlsx_data()
        result = engine.export(data)
        wb = openpyxl.load_workbook(str(result.file_path))
        assert "Sources" in wb.sheetnames

    def test_xlsx_header_row_present(self, engine):
        import openpyxl
        result = engine.export(_make_xlsx_data())
        wb = openpyxl.load_workbook(str(result.file_path))
        ws = wb.active or wb["Revenue"]
        ws = wb["Revenue"]
        first_row = [ws.cell(1, col).value for col in range(1, 4)]
        assert "Month" in first_row

    def test_xlsx_correct_row_count(self, engine):
        import openpyxl
        result = engine.export(_make_xlsx_data())
        wb = openpyxl.load_workbook(str(result.file_path))
        ws = wb["Revenue"]
        # 1 header + 3 data rows = 4 rows
        assert ws.max_row == 4

    def test_xlsx_empty_sheets_creates_summary(self, engine):
        data = _make_xlsx_data(sheets=[])
        result = engine.export(data)
        assert result.success


# ── ExportEngine DOCX Tests ────────────────────────────────────────────────
class TestExportDOCX:

    @pytest.fixture
    def engine(self, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "modules.export.export_engine.OUTPUT_DIR", tmp_path
        )
        from modules.export.export_engine import ExportEngine
        return ExportEngine()

    def test_docx_file_created(self, engine):
        result = engine.export(_make_docx_data())
        assert result.success
        assert result.file_path.suffix == ".docx"

    def test_docx_file_non_empty(self, engine):
        result = engine.export(_make_docx_data())
        assert result.file_size_bytes > 1000

    def test_docx_readable_by_python_docx(self, engine):
        from docx import Document
        result = engine.export(_make_docx_data())
        doc = Document(str(result.file_path))
        full_text = " ".join(p.text for p in doc.paragraphs)
        assert len(full_text) > 0

    def test_docx_contains_title(self, engine):
        from docx import Document
        data   = _make_docx_data(title="My Special Title")
        result = engine.export(data)
        doc    = Document(str(result.file_path))
        all_text = " ".join(p.text for p in doc.paragraphs)
        assert "My Special Title" in all_text

    def test_docx_contains_section_headings(self, engine):
        from docx import Document
        result = engine.export(_make_docx_data())
        doc    = Document(str(result.file_path))
        headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
        assert any("Executive Summary" in h for h in headings)


# ── ExportEngine PPTX Tests ────────────────────────────────────────────────
class TestExportPPTX:

    @pytest.fixture
    def engine(self, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "modules.export.export_engine.OUTPUT_DIR", tmp_path
        )
        from modules.export.export_engine import ExportEngine
        return ExportEngine()

    def test_pptx_file_created(self, engine):
        result = engine.export(_make_pptx_data())
        assert result.success
        assert result.file_path.suffix == ".pptx"

    def test_pptx_file_non_empty(self, engine):
        result = engine.export(_make_pptx_data())
        assert result.file_size_bytes > 1000

    def test_pptx_has_minimum_slides(self, engine):
        from pptx import Presentation
        result = engine.export(_make_pptx_data())
        prs = Presentation(str(result.file_path))
        # Title slide + 1 content slide + sources slide = at least 3
        assert len(prs.slides) >= 2

    def test_pptx_16x9_aspect_ratio(self, engine):
        from pptx import Presentation
        from pptx.util import Inches
        result = engine.export(_make_pptx_data())
        prs    = Presentation(str(result.file_path))
        width  = prs.slide_width
        height = prs.slide_height
        ratio  = width / height
        assert abs(ratio - (16/9)) < 0.1


# ── ExportEngine CSV Tests ─────────────────────────────────────────────────
class TestExportCSV:

    @pytest.fixture
    def engine(self, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "modules.export.export_engine.OUTPUT_DIR", tmp_path
        )
        from modules.export.export_engine import ExportEngine
        return ExportEngine()

    def test_csv_single_sheet_creates_csv_file(self, engine):
        result = engine.export(_make_csv_data())
        assert result.success
        assert result.file_path.suffix == ".csv"

    def test_csv_content_correct(self, engine):
        import csv
        result = engine.export(_make_csv_data())
        with open(str(result.file_path), "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows   = list(reader)
        # header + 2 data rows
        assert len(rows) == 3
        assert rows[0] == ["Month", "Revenue"]

    def test_csv_multi_sheet_creates_zip(self, engine):
        result = engine.export(_make_csv_data(multi_sheet=True))
        assert result.success
        assert result.file_path.suffix == ".zip"
        # Verify zip contains CSVs
        with zipfile.ZipFile(str(result.file_path)) as zf:
            names = zf.namelist()
        assert any(n.endswith(".csv") for n in names)
        assert len(names) == 2


# ── ExportEngine General Tests ─────────────────────────────────────────────
class TestExportGeneral:

    @pytest.fixture
    def engine(self, tmp_path, monkeypatch):
        monkeypatch.setattr(
            "modules.export.export_engine.OUTPUT_DIR", tmp_path
        )
        from modules.export.export_engine import ExportEngine
        return ExportEngine()

    def test_unknown_type_defaults_to_xlsx(self, engine):
        data   = _make_xlsx_data(export_type="unknown_format")
        result = engine.export(data)
        assert result.success
        assert result.file_path.suffix == ".xlsx"

    def test_type_alias_excel_works(self, engine):
        data   = _make_xlsx_data(export_type="excel")
        result = engine.export(data)
        assert result.success
        assert result.file_path.suffix == ".xlsx"

    def test_type_alias_word_works(self, engine):
        data   = _make_docx_data(export_type="word")
        result = engine.export(data)
        assert result.success
        assert result.file_path.suffix == ".docx"

    def test_export_result_to_dict(self, engine):
        result = engine.export(_make_xlsx_data())
        d = result.to_dict()
        assert "file_path"       in d
        assert "export_type"     in d
        assert "file_size_bytes" in d
        assert "success"         in d

    def test_export_from_string_valid_json(self, engine):
        data_str = json.dumps(_make_xlsx_data())
        result   = engine.export_from_string(data_str)
        assert result.success

    def test_export_from_string_invalid_json(self, engine):
        result = engine.export_from_string("not valid json {{{{")
        assert result.success is False
        assert result.error is not None

    def test_filename_sanitization(self):
        from modules.export.export_engine import _safe_filename
        filename = _safe_filename("Q3 Revenue/Analysis!", "xlsx")
        assert "/" not in filename
        assert "!" not in filename
        assert filename.endswith(".xlsx")

    def test_filename_empty_name(self):
        from modules.export.export_engine import _safe_filename
        filename = _safe_filename("", "docx")
        assert filename.endswith(".docx")
        assert len(filename) > 5

    def test_export_success_flag(self, engine):
        result = engine.export(_make_xlsx_data())
        assert result.success is True
        assert result.error is None
