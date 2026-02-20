"""
modules/export/export_engine.py

File Export Engine for the Multimodal RAG system.

Converts the structured JSON dict produced by LLMEngine.export()
into real files: DOCX, XLSX, PPTX, or CSV.

Export types:
  xlsx — Excel workbook (openpyxl)
         Professional styling: header row, alternating rows,
         auto column width, freeze top row, sources sheet
  docx — Word document (python-docx)
         Title, summary, sections with headings,
         sources table, page numbers in footer
  pptx — PowerPoint presentation (python-pptx)
         Title slide, content slides, sources slide
         Professional 16:9 layout
  csv  — CSV file(s) (stdlib csv module)
         One file per sheet, zipped if multiple

All files saved to OUTPUT_DIR with timestamped filenames.
"""

import csv
import io
import json
import re
import zipfile
from datetime import datetime, timezone
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from loguru import logger

from config.settings import OUTPUT_DIR


# ── Export Result ──────────────────────────────────────────────────────────
@dataclass
class ExportResult:
    """
    Result of a file export operation.
    Returned to the API layer which serves the file to the user.
    """
    file_path:       Path
    export_type:     str
    filename:        str
    file_size_bytes: int
    created_at:      str
    success:         bool
    error:           Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_path":       str(self.file_path),
            "export_type":     self.export_type,
            "filename":        self.filename,
            "file_size_bytes": self.file_size_bytes,
            "created_at":      self.created_at,
            "success":         self.success,
            "error":           self.error,
        }


# ── Filename Utilities ─────────────────────────────────────────────────────
def _safe_filename(name: str, extension: str) -> str:
    """
    Sanitize a filename and append timestamp + extension.
    Removes special characters, collapses spaces, limits length.

    Example: "Q3 Revenue/Analysis!" → "Q3_Revenue_Analysis_20240315_143022.xlsx"
    """
    # Remove unsafe characters
    safe = re.sub(r"[^\w\s\-]", "", name)
    safe = re.sub(r"\s+", "_", safe.strip())
    safe = safe[:50] if len(safe) > 50 else safe
    safe = safe or "export"

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{safe}_{timestamp}.{extension}"


def _get_output_path(filename: str) -> Path:
    """Ensure OUTPUT_DIR exists and return full path for filename."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR / filename


# ── XLSX Export ────────────────────────────────────────────────────────────
def _export_xlsx(export_data: Dict[str, Any]) -> Path:
    """
    Create a professionally styled Excel workbook from export data.

    Styling:
    - Header row: dark navy background (#1E3A5F), white bold text, 12pt
    - Data rows: alternating white (#FFFFFF) and light blue (#EBF5FB)
    - All columns auto-width (max 50 chars)
    - Top row frozen for scrolling
    - Gridlines visible
    - Sources appended as separate sheet
    - Workbook title set from export_data["title"]
    """
    try:
        import openpyxl
        from openpyxl.styles import (
            Font, PatternFill, Alignment, Border, Side
        )
        from openpyxl.utils import get_column_letter
    except ImportError as e:
        raise RuntimeError(f"openpyxl not installed: {e}")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default empty sheet

    # ── Color constants ────────────────────────────────────────────────────
    HEADER_FILL  = PatternFill("solid", fgColor="1E3A5F")
    ALT_ROW_FILL = PatternFill("solid", fgColor="EBF5FB")
    WHITE_FILL   = PatternFill("solid", fgColor="FFFFFF")

    HEADER_FONT  = Font(name="Calibri", bold=True,  color="FFFFFF", size=12)
    DATA_FONT    = Font(name="Calibri", bold=False, color="000000", size=11)

    CENTER_ALIGN = Alignment(horizontal="center", vertical="center", wrap_text=False)
    LEFT_ALIGN   = Alignment(horizontal="left",   vertical="center", wrap_text=True)

    THIN_BORDER  = Border(
        left=Side(style="thin"),  right=Side(style="thin"),
        top=Side(style="thin"),   bottom=Side(style="thin"),
    )

    sheets = export_data.get("sheets", [])
    if not sheets:
        # Create a minimal sheet with the summary
        sheets = [{
            "name": "Summary",
            "headers": ["Field", "Value"],
            "rows": [
                ["Title",    export_data.get("title",   "")],
                ["Summary",  export_data.get("summary", "")],
            ],
        }]

    for sheet_data in sheets:
        sheet_name = str(sheet_data.get("name", "Sheet"))[:31]  # Excel limit
        headers    = sheet_data.get("headers", [])
        rows       = sheet_data.get("rows", [])

        ws = wb.create_sheet(title=sheet_name)

        # ── Write header row ───────────────────────────────────────────────
        if headers:
            for col_idx, header in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=col_idx, value=str(header))
                cell.font      = HEADER_FONT
                cell.fill      = HEADER_FILL
                cell.alignment = CENTER_ALIGN
                cell.border    = THIN_BORDER
            ws.row_dimensions[1].height = 22

        # ── Write data rows ────────────────────────────────────────────────
        for row_idx, row in enumerate(rows, start=2):
            fill = ALT_ROW_FILL if row_idx % 2 == 0 else WHITE_FILL
            for col_idx, value in enumerate(row, start=1):
                # Convert value to appropriate type
                cell_value = value
                if isinstance(value, (int, float)):
                    cell_value = value
                elif isinstance(value, str):
                    # Try to parse numeric strings
                    try:
                        if "." in value:
                            cell_value = float(value)
                        else:
                            cell_value = int(value)
                    except (ValueError, TypeError):
                        cell_value = str(value)

                cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                cell.font      = DATA_FONT
                cell.fill      = fill
                cell.alignment = LEFT_ALIGN
                cell.border    = THIN_BORDER

        # ── Auto column width ──────────────────────────────────────────────
        for col_idx, col_cells in enumerate(ws.columns, start=1):
            max_length = 0
            col_letter = get_column_letter(col_idx)
            for cell in col_cells:
                try:
                    cell_len = len(str(cell.value or ""))
                    if cell_len > max_length:
                        max_length = cell_len
                except Exception:
                    pass
            adjusted = min(max_length + 4, 50)
            ws.column_dimensions[col_letter].width = adjusted

        # ── Freeze top row ─────────────────────────────────────────────────
        ws.freeze_panes = "A2"

    # ── Sources sheet ──────────────────────────────────────────────────────
    sources = export_data.get("sources", [])
    if sources:
        ws_src = wb.create_sheet(title="Sources")
        ws_src.cell(row=1, column=1, value="Source File").font = HEADER_FONT
        ws_src.cell(row=1, column=1).fill = HEADER_FILL
        ws_src.cell(row=1, column=2, value="Reference").font = HEADER_FONT
        ws_src.cell(row=1, column=2).fill = HEADER_FILL

        for i, src in enumerate(sources, start=2):
            parts = str(src).split("|", 1)
            ws_src.cell(row=i, column=1, value=parts[0].strip())
            ws_src.cell(row=i, column=2, value=parts[1].strip() if len(parts) > 1 else "")

        ws_src.column_dimensions["A"].width = 40
        ws_src.column_dimensions["B"].width = 20
        ws_src.freeze_panes = "A2"

    # ── Save ───────────────────────────────────────────────────────────────
    base_name = export_data.get("filename", "export")
    filename  = _safe_filename(base_name, "xlsx")
    out_path  = _get_output_path(filename)
    wb.save(str(out_path))

    logger.info(f"XLSX exported: {out_path.name} ({out_path.stat().st_size} bytes)")
    return out_path


# ── DOCX Export ────────────────────────────────────────────────────────────
def _export_docx(export_data: Dict[str, Any]) -> Path:
    """
    Create a professionally styled Word document from export data.

    Structure:
    - Title (Heading 1) from export_data["title"]
    - Summary paragraph (italic, indented)
    - Section headings (Heading 2) and body paragraphs from "sections"
    - Horizontal rule before Sources
    - Sources table (2 cols: Source File | Reference)
    - Page numbers in footer (right-aligned)
    - 1.15 line spacing throughout
    """
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError as e:
        raise RuntimeError(f"python-docx not installed: {e}")

    doc = Document()

    # ── Page margins ───────────────────────────────────────────────────────
    section = doc.sections[0]
    section.top_margin    = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin   = Cm(3.0)
    section.right_margin  = Cm(3.0)

    # ── Title ──────────────────────────────────────────────────────────────
    title = export_data.get("title", "Document")
    title_para = doc.add_heading(title, level=1)
    title_para.runs[0].font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    # ── Summary ────────────────────────────────────────────────────────────
    summary = export_data.get("summary", "")
    if summary:
        summary_para = doc.add_paragraph()
        summary_run  = summary_para.add_run(summary)
        summary_run.italic = True
        summary_run.font.size = Pt(11)
        summary_run.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
        summary_para.paragraph_format.space_after = Pt(12)

    doc.add_paragraph()  # spacing

    # ── Sections ───────────────────────────────────────────────────────────
    sections = export_data.get("sections", [])
    if not sections:
        # Fall back: use sheets as sections if no sections key
        for sheet in export_data.get("sheets", []):
            sections.append({
                "heading": sheet.get("name", "Data"),
                "content": _sheet_to_text(sheet),
            })

    for sec in sections:
        heading = sec.get("heading", "")
        content = sec.get("content", "")

        if heading:
            h2 = doc.add_heading(heading, level=2)
            if h2.runs:
                h2.runs[0].font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)

        if content:
            para = doc.add_paragraph(content)
            para.paragraph_format.space_after = Pt(8)
            if para.runs:
                para.runs[0].font.size = Pt(11)

    # ── Sources Table ──────────────────────────────────────────────────────
    sources = export_data.get("sources", [])
    if sources:
        doc.add_paragraph()
        doc.add_heading("Sources", level=2)

        table = doc.add_table(rows=1, cols=2)
        table.style = "Table Grid"

        # Header row
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = "Source File"
        hdr_cells[1].text = "Reference"
        for cell in hdr_cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Set header cell background
            _set_cell_background(cell, "1E3A5F")

        # Data rows
        for src in sources:
            parts = str(src).split("|", 1)
            row_cells = table.add_row().cells
            row_cells[0].text = parts[0].strip()
            row_cells[1].text = parts[1].strip() if len(parts) > 1 else ""

        # Column widths
        for row in table.rows:
            row.cells[0].width = Inches(3.5)
            row.cells[1].width = Inches(2.0)

    # ── Footer with page numbers ───────────────────────────────────────────
    _add_page_numbers_footer(doc)

    # ── Save ───────────────────────────────────────────────────────────────
    base_name = export_data.get("filename", "document")
    filename  = _safe_filename(base_name, "docx")
    out_path  = _get_output_path(filename)
    doc.save(str(out_path))

    logger.info(f"DOCX exported: {out_path.name} ({out_path.stat().st_size} bytes)")
    return out_path


def _sheet_to_text(sheet: Dict[str, Any]) -> str:
    """Convert a sheet dict to plain text for DOCX sections."""
    headers = sheet.get("headers", [])
    rows    = sheet.get("rows", [])
    lines   = []
    if headers:
        lines.append(" | ".join(str(h) for h in headers))
        lines.append("-" * 40)
    for row in rows:
        lines.append(" | ".join(str(v) for v in row))
    return "\n".join(lines)


def _set_cell_background(cell, hex_color: str):
    """Set DOCX table cell background color using XML manipulation."""
    try:
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"),   "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"),  hex_color)
        tcPr.append(shd)
    except Exception as e:
        logger.debug(f"Could not set cell background: {e}")


def _add_page_numbers_footer(doc):
    """Add a right-aligned page number to the DOCX footer."""
    try:
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        section = doc.sections[0]
        footer  = section.footer
        para    = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()

        from docx.enum.text import WD_ALIGN_PARAGRAPH
        para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        run = para.add_run()
        fldChar1 = OxmlElement("w:fldChar")
        fldChar1.set(qn("w:fldCharType"), "begin")

        instrText = OxmlElement("w:instrText")
        instrText.text = "PAGE"

        fldChar2 = OxmlElement("w:fldChar")
        fldChar2.set(qn("w:fldCharType"), "end")

        run._r.append(fldChar1)
        run._r.append(instrText)
        run._r.append(fldChar2)
    except Exception as e:
        logger.debug(f"Could not add page numbers: {e}")


# ── PPTX Export ────────────────────────────────────────────────────────────
def _export_pptx(export_data: Dict[str, Any]) -> Path:
    """
    Create a professionally styled PowerPoint presentation.

    Slide structure:
    - Slide 1: Title slide (title + summary as subtitle)
    - Slides 2..N: Content slides, one per section
                   Title from section heading
                   Body from bullet points or content text
    - Final slide: Sources slide
                   Title "Sources" + bulleted source list

    Layout: 16:9 widescreen (13.33" × 7.5")
    Theme: Navy/white professional
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        raise RuntimeError(f"python-pptx not installed: {e}")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    BLUE   = RGBColor(0x1E, 0x40, 0xAF)
    LGRAY  = RGBColor(0xF3, 0xF4, 0xF6)

    blank_layout = prs.slide_layouts[6]   # completely blank

    def _add_text_box(slide, text, left, top, width, height,
                      font_size=24, bold=False, color=None,
                      align=PP_ALIGN.LEFT, italic=False):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txBox

    def _add_background(slide, color_rgb: RGBColor):
        """Fill slide background with a solid color."""
        from pptx.util import Emu
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    # ── Slide 1: Title slide ───────────────────────────────────────────────
    title_slide = prs.slides.add_slide(blank_layout)
    _add_background(title_slide, NAVY)

    title_text   = export_data.get("title", "Report")
    summary_text = export_data.get("summary", "")

    _add_text_box(
        title_slide, title_text,
        left=0.8, top=2.0, width=11.7, height=1.5,
        font_size=40, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT,
    )

    if summary_text:
        summary_truncated = summary_text[:300]
        _add_text_box(
            title_slide, summary_truncated,
            left=0.8, top=3.8, width=11.7, height=1.5,
            font_size=20, bold=False, color=RGBColor(0xBF, 0xD7, 0xFF),
            italic=True, align=PP_ALIGN.LEFT,
        )

    # Thin accent line at bottom
    from pptx.util import Inches as In
    line = title_slide.shapes.add_connector(
        1,   # MSO_CONNECTOR_TYPE.STRAIGHT
        In(0.8), In(6.8), In(12.5), In(6.8)
    )
    line.line.color.rgb = RGBColor(0x60, 0xA5, 0xFA)
    line.line.width     = Pt(2)

    # ── Content slides ─────────────────────────────────────────────────────
    sections = export_data.get("sections", [])
    if not sections:
        for sheet in export_data.get("sheets", []):
            bullet_pts = []
            headers = sheet.get("headers", [])
            for row in sheet.get("rows", [])[:10]:  # max 10 rows per slide
                bullet_pts.append(" | ".join(str(v) for v in row))
            sections.append({
                "heading":       sheet.get("name", "Data"),
                "bullet_points": bullet_pts,
                "content":       "",
            })

    for sec in sections:
        slide = prs.slides.add_slide(blank_layout)
        _add_background(slide, WHITE)

        # Navy top bar
        bar = slide.shapes.add_shape(
            1,   # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()   # no border

        # Section title
        heading = sec.get("heading", "")
        _add_text_box(
            slide, heading,
            left=0.4, top=0.15, width=12.5, height=0.9,
            font_size=28, bold=True, color=WHITE,
        )

        # Bullet points or content
        bullet_points = sec.get("bullet_points", [])
        content       = sec.get("content", "")

        if bullet_points:
            # Add bullet point text box
            txBox = slide.shapes.add_textbox(
                Inches(0.6), Inches(1.5),
                Inches(12.1), Inches(5.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, bp in enumerate(bullet_points[:10]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text  = f"• {str(bp)[:200]}"
                p.level = 0
                if p.runs:
                    p.runs[0].font.size  = Pt(18)
                    p.runs[0].font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

        elif content:
            _add_text_box(
                slide, content[:600],
                left=0.6, top=1.5, width=12.1, height=5.5,
                font_size=16, color=RGBColor(0x1F, 0x29, 0x37),
            )

    # ── Sources slide ──────────────────────────────────────────────────────
    sources = export_data.get("sources", [])
    if sources:
        src_slide = prs.slides.add_slide(blank_layout)
        _add_background(src_slide, LGRAY)

        bar2 = src_slide.shapes.add_shape(
            1,
            Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = NAVY
        bar2.line.fill.background()

        _add_text_box(
            src_slide, "Sources",
            left=0.4, top=0.15, width=12.5, height=0.9,
            font_size=28, bold=True, color=WHITE,
        )

        src_text = "\n".join(
            f"• {str(s)}" for s in sources[:15]
        )
        _add_text_box(
            src_slide, src_text,
            left=0.6, top=1.5, width=12.1, height=5.5,
            font_size=16, color=RGBColor(0x1F, 0x29, 0x37),
        )

    # ── Save ───────────────────────────────────────────────────────────────
    base_name = export_data.get("filename", "presentation")
    filename  = _safe_filename(base_name, "pptx")
    out_path  = _get_output_path(filename)
    prs.save(str(out_path))

    logger.info(f"PPTX exported: {out_path.name} ({out_path.stat().st_size} bytes)")
    return out_path


# ── CSV Export ─────────────────────────────────────────────────────────────
def _export_csv(export_data: Dict[str, Any]) -> Path:
    """
    Create CSV file(s) from export data.
    If there is only one sheet, writes a single .csv file.
    If there are multiple sheets, writes multiple CSVs and zips them.
    """
    sheets = export_data.get("sheets", [])
    if not sheets:
        sheets = [{
            "name": "data",
            "headers": ["Summary"],
            "rows": [[export_data.get("summary", "")]],
        }]

    base_name = export_data.get("filename", "export")

    if len(sheets) == 1:
        filename = _safe_filename(base_name, "csv")
        out_path = _get_output_path(filename)
        with open(str(out_path), "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            headers = sheets[0].get("headers", [])
            if headers:
                writer.writerow(headers)
            for row in sheets[0].get("rows", []):
                writer.writerow([str(v) for v in row])
        logger.info(f"CSV exported: {out_path.name}")
        return out_path

    else:
        # Multiple sheets → zip
        zip_filename = _safe_filename(base_name, "zip")
        zip_path     = _get_output_path(zip_filename)
        with zipfile.ZipFile(str(zip_path), "w", zipfile.ZIP_DEFLATED) as zf:
            for i, sheet in enumerate(sheets):
                sheet_name = re.sub(r"[^\w\-]", "_", sheet.get("name", f"sheet_{i}"))
                csv_name   = f"{sheet_name}.csv"
                buf = io.StringIO()
                writer = csv.writer(buf)
                headers = sheet.get("headers", [])
                if headers:
                    writer.writerow(headers)
                for row in sheet.get("rows", []):
                    writer.writerow([str(v) for v in row])
                zf.writestr(csv_name, buf.getvalue())

        logger.info(f"CSV ZIP exported: {zip_path.name}")
        return zip_path


# ── Public API ─────────────────────────────────────────────────────────────
class ExportEngine:
    """
    Converts LLM export JSON into real files (DOCX, XLSX, PPTX, CSV).

    Usage:
        engine = ExportEngine()

        # After LLMEngine.export() returns a JSON dict:
        export_data = llm_engine.export(query, retrieval_result)
        result = engine.export(export_data)

        if result.success:
            # Serve result.file_path to the user
            print(f"File ready: {result.file_path}")
        else:
            print(f"Export failed: {result.error}")
    """

    def export(self, export_data: Dict[str, Any]) -> ExportResult:
        """
        Convert export JSON data to a file on disk.

        Args:
            export_data: Dict from LLMEngine.export() with keys:
                         export_type, filename, title, summary,
                         sheets (xlsx/csv), sections (docx),
                         slides (pptx), sources

        Returns:
            ExportResult with file_path and metadata.
        """
        export_type = str(
            export_data.get("export_type", "xlsx")
        ).lower().strip()

        now = datetime.now(timezone.utc).isoformat()

        # Normalize export type aliases
        type_aliases = {
            "excel": "xlsx", "spreadsheet": "xlsx",
            "word":  "docx", "document":    "docx",
            "powerpoint": "pptx", "presentation": "pptx", "slides": "pptx",
        }
        export_type = type_aliases.get(export_type, export_type)

        if export_type not in ("xlsx", "docx", "pptx", "csv"):
            logger.warning(
                f"Unknown export type '{export_type}' — defaulting to xlsx"
            )
            export_type = "xlsx"

        logger.info(
            f"ExportEngine: generating {export_type.upper()} — "
            f"'{export_data.get('title', 'untitled')}'"
        )

        try:
            if export_type == "xlsx":
                file_path = _export_xlsx(export_data)
            elif export_type == "docx":
                file_path = _export_docx(export_data)
            elif export_type == "pptx":
                file_path = _export_pptx(export_data)
            elif export_type == "csv":
                file_path = _export_csv(export_data)
            else:
                file_path = _export_xlsx(export_data)

            return ExportResult(
                file_path=file_path,
                export_type=export_type,
                filename=file_path.name,
                file_size_bytes=file_path.stat().st_size,
                created_at=now,
                success=True,
            )

        except Exception as e:
            logger.error(f"Export failed: {e}")
            return ExportResult(
                file_path=Path(""),
                export_type=export_type,
                filename="",
                file_size_bytes=0,
                created_at=now,
                success=False,
                error=str(e),
            )

    def export_from_string(self, json_string: str) -> ExportResult:
        """
        Parse a JSON string and export.
        Convenience wrapper for when the LLM output is not yet parsed.
        """
        try:
            export_data = json.loads(json_string)
        except json.JSONDecodeError as e:
            now = datetime.now(timezone.utc).isoformat()
            return ExportResult(
                file_path=Path(""),
                export_type="unknown",
                filename="",
                file_size_bytes=0,
                created_at=now,
                success=False,
                error=f"Invalid JSON: {e}",
            )
        return self.export(export_data)
